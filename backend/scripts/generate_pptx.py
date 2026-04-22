#!/usr/bin/env python3
"""PPTX Generator for ConvoiaAI. JSON in via stdin, PPTX written to argv[1]."""
import sys
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


def main():
    output_path = sys.argv[1]
    data = json.load(sys.stdin)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    theme = data.get('theme', {})
    primary = RGBColor.from_string(theme.get('primaryColor', '#7C3AED').lstrip('#'))

    for slide_data in data.get('slides', []):
        stype = slide_data.get('type', 'content')

        if stype == 'title':
            layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(layout)
            if slide.shapes.title:
                slide.shapes.title.text = slide_data.get('title', '')
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = slide_data.get('subtitle', '')

        elif stype == 'content':
            layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = slide_data.get('title', '')
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()
            bullets = slide_data.get('bullets', [])
            for i, bullet in enumerate(bullets):
                if i == 0:
                    tf.text = bullet
                else:
                    p = tf.add_paragraph()
                    p.text = bullet
                    p.level = 0
            if slide_data.get('notes'):
                slide.notes_slide.notes_text_frame.text = slide_data['notes']

        elif stype == 'two_column':
            layout = prs.slide_layouts[3]
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = slide_data.get('title', '')

            if len(slide.placeholders) > 1:
                tf_left = slide.placeholders[1].text_frame
                tf_left.clear()
                if slide_data.get('left_title'):
                    tf_left.text = slide_data['left_title']
                    tf_left.paragraphs[0].font.bold = True
                for bullet in slide_data.get('left_bullets', []):
                    p = tf_left.add_paragraph()
                    p.text = bullet

            if len(slide.placeholders) > 2:
                tf_right = slide.placeholders[2].text_frame
                tf_right.clear()
                if slide_data.get('right_title'):
                    tf_right.text = slide_data['right_title']
                    tf_right.paragraphs[0].font.bold = True
                for bullet in slide_data.get('right_bullets', []):
                    p = tf_right.add_paragraph()
                    p.text = bullet

        elif stype == 'table':
            layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(layout)

            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
            tf = txBox.text_frame
            tf.text = slide_data.get('title', '')
            tf.paragraphs[0].font.size = Pt(28)
            tf.paragraphs[0].font.bold = True

            headers = slide_data.get('headers', [])
            rows = slide_data.get('rows', [])
            cols = max(len(headers), 1)
            total_rows = 1 + len(rows)

            table_shape = slide.shapes.add_table(
                total_rows, cols,
                Inches(0.5), Inches(1.5),
                Inches(12), Inches(0.4 * total_rows),
            )
            table = table_shape.table
            for i, header in enumerate(headers):
                cell = table.cell(0, i)
                cell.text = str(header)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.bold = True
                    paragraph.font.size = Pt(12)
                    paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                cell.fill.solid()
                cell.fill.fore_color.rgb = primary
            for r, row in enumerate(rows):
                for c, value in enumerate(row):
                    if c < cols:
                        table.cell(r + 1, c).text = str(value)

    prs.save(output_path)
    print(json.dumps({'success': True, 'path': output_path}))


if __name__ == '__main__':
    main()
